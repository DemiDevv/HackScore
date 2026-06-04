import asyncio
import re
from datetime import UTC, datetime
from pathlib import Path
from uuid import UUID

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PyPDF2 import PdfReader
from sqlalchemy import select

from app.database import AsyncSessionLocal
from app.models.check_result import CheckResult
from app.models.enums import CheckStatus, CheckType
from app.models.submission import Submission
from app.workers.celery_app import celery_app


SECTION_PATTERNS = {
    "problem": re.compile(r"(проблема|problem)", re.IGNORECASE),
    "solution": re.compile(r"(решение|solution)", re.IGNORECASE),
    "target_audience": re.compile(r"(целевая аудитория|target audience|\bца\b)", re.IGNORECASE),
    "stack": re.compile(r"(стек|технологи|stack|tech)", re.IGNORECASE),
    "demo": re.compile(r"(демо|demo)", re.IGNORECASE),
    "team": re.compile(r"(команда|team)", re.IGNORECASE),
}


def slide_count_score(count: int) -> float:
    if count == 0:
        return 0.0
    if count <= 4:
        return 3.0
    if count <= 7:
        return 7.0
    if count <= 15:
        return 10.0
    if count <= 25:
        return 7.0
    return 5.0


def extract_shape_text(shape: object) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs)


def read_pptx(path: Path) -> tuple[int, str, bool, bool]:
    presentation = Presentation(str(path))
    text_parts: list[str] = []
    has_images = False
    has_diagrams = False

    for slide in presentation.slides:
        for shape in slide.shapes:
            text_parts.append(extract_shape_text(shape))
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_images = True
            if getattr(shape, "has_chart", False) or getattr(shape, "has_table", False) or shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                has_diagrams = True

    return len(presentation.slides), "\n".join(text_parts), has_images, has_diagrams


def read_pdf(path: Path) -> tuple[int, str, bool, bool]:
    reader = PdfReader(str(path))
    text = "\n".join(page.extract_text() or "" for page in reader.pages)
    has_images = False

    for page in reader.pages:
        resources = page.get("/Resources") or {}
        x_objects = resources.get("/XObject") or {}
        try:
            objects = x_objects.get_object()
        except Exception:
            objects = {}
        for obj in objects.values():
            try:
                if obj.get_object().get("/Subtype") == "/Image":
                    has_images = True
                    break
            except Exception:
                continue
        if has_images:
            break

    return len(reader.pages), text, has_images, False


def extract_presentation(path: Path) -> tuple[int, str, bool, bool]:
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        return read_pptx(path)
    if suffix == ".pdf":
        return read_pdf(path)
    raise RuntimeError("Unsupported presentation format")


def analyze_presentation(path: Path) -> tuple[float, dict[str, object]]:
    slide_count, text, has_images, has_diagrams = extract_presentation(path)
    sections = {name: bool(pattern.search(text)) for name, pattern in SECTION_PATTERNS.items()}
    sections_found = [name for name, found in sections.items() if found]
    sections_missing = [name for name, found in sections.items() if not found]

    count_score = slide_count_score(slide_count)
    structure_score = round((len(sections_found) / len(SECTION_PATTERNS)) * 10, 2)
    visual_bonus = (1.0 if has_images else 0.0) + (1.0 if has_diagrams else 0.0)
    total = min(10.0, count_score * 0.45 + structure_score * 0.45 + visual_bonus)

    report = {
        "slide_count": slide_count,
        "sections_found": sections_found,
        "sections_missing": sections_missing,
        "has_visuals": has_images or has_diagrams,
        "has_images": has_images,
        "has_diagrams": has_diagrams,
        "scores": {
            "slide_count": count_score,
            "structure": structure_score,
            "visual_bonus": visual_bonus,
        },
    }
    return round(total, 2), report


async def get_presentation_check(submission_id: UUID) -> tuple[Submission, CheckResult]:
    async with AsyncSessionLocal() as db:
        submission = await db.get(Submission, submission_id)
        if submission is None:
            raise RuntimeError("Submission not found")

        result = await db.execute(
            select(CheckResult).where(
                CheckResult.submission_id == submission_id,
                CheckResult.check_type == CheckType.presentation,
            )
        )
        check = result.scalar_one_or_none()
        if check is None:
            check = CheckResult(
                submission_id=submission_id,
                check_type=CheckType.presentation,
                status=CheckStatus.pending,
            )
            db.add(check)
            await db.commit()
            await db.refresh(check)
        return submission, check


async def mark_running(check_id: UUID) -> None:
    async with AsyncSessionLocal() as db:
        check = await db.get(CheckResult, check_id)
        if check is None:
            return
        check.status = CheckStatus.running
        check.started_at = datetime.now(UTC)
        await db.commit()


async def save_check_result(check_id: UUID, status_value: CheckStatus, score: float | None, report: dict[str, object]) -> None:
    async with AsyncSessionLocal() as db:
        check = await db.get(CheckResult, check_id)
        if check is None:
            return
        check.status = status_value
        check.score = score
        check.report = report
        check.completed_at = datetime.now(UTC)
        await db.commit()


async def run_check(submission_id: str) -> None:
    submission, check = await get_presentation_check(UUID(submission_id))
    await mark_running(check.id)

    try:
        if not submission.presentation:
            raise RuntimeError("Presentation file is missing")
        score, report = analyze_presentation(Path(submission.presentation))
        await save_check_result(check.id, CheckStatus.completed, score, report)
    except Exception as exc:
        await save_check_result(check.id, CheckStatus.failed, None, {"error": str(exc)})


@celery_app.task(name="app.workers.presentation_checker.check_presentation")
def check_presentation(submission_id: str) -> None:
    asyncio.run(run_check(submission_id))
